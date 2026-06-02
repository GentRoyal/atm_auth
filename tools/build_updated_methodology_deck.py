from __future__ import annotations

import math
from pathlib import Path

import fitz
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE_PDF = Path(
    r"C:\Users\USER\Desktop\Design_and_Implementation_of_a_Multi-Factor_Biometric_Authentication_System_for_ATM_Security.pptx.pdf"
)
OUT_DIR = ROOT / "outputs"
PAGE_DIR = OUT_DIR / "pdf_pages"
ASSET_DIR = OUT_DIR / "assets"
PREVIEW_DIR = OUT_DIR / "previews"
PPTX_PATH = OUT_DIR / "ATM_Biometric_Updated_Methodology.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(79, 129, 189)
DARK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(20, 137, 64)
RED = RGBColor(193, 23, 24)
YELLOW = RGBColor(255, 223, 89)
NAVY = RGBColor(28, 33, 59)


def ensure_dirs() -> None:
    for path in (OUT_DIR, PAGE_DIR, ASSET_DIR, PREVIEW_DIR):
        path.mkdir(parents=True, exist_ok=True)


def render_pdf_pages() -> list[Path]:
    doc = fitz.open(SOURCE_PDF)
    paths: list[Path] = []
    matrix = fitz.Matrix(2, 2)
    for idx, page in enumerate(doc, start=1):
        path = PAGE_DIR / f"page_{idx:02d}.png"
        page.get_pixmap(matrix=matrix, alpha=False).save(path)
        paths.append(path)
    return paths


def crop_assets(page_paths: list[Path]) -> dict[str, Path]:
    assets: dict[str, Path] = {}
    page7 = Image.open(page_paths[6]).convert("RGB")
    w, h = page7.size

    crops = {
        "footer": (0, int(h * 0.91), w, h),
        "chevron": (0, 0, int(w * 0.12), int(h * 0.19)),
        "integration": (int(w * 0.735), int(h * 0.12), int(w * 0.955), int(h * 0.50)),
    }
    for name, box in crops.items():
        out = ASSET_DIR / f"{name}.png"
        page7.crop(box).save(out)
        assets[name] = out
    return assets


def add_picture_slide(prs: Presentation, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_shared_branding(slide, assets: dict[str, Path]) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_picture(str(assets["chevron"]), 0, 0, width=Inches(1.55), height=Inches(1.42))
    slide.shapes.add_picture(str(assets["footer"]), 0, Inches(6.82), width=SLIDE_W, height=Inches(0.68))


def add_textbox(slide, text, x, y, w, h, size=24, bold=False, italic=False, color=DARK,
                font="Times New Roman", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_stage_line(slide, num, title, detail, y):
    add_textbox(slide, f"{num}.", 2.18, y, 0.35, 0.35, size=22)
    add_textbox(slide, title, 2.55, y, 3.0, 0.35, size=22, bold=True)
    add_textbox(slide, f"- {detail}", 5.28, y, 3.55, 0.35, size=20)


def add_methodology_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shared_branding(slide, assets)

    add_textbox(slide, "METHODOLOGY", 2.18, 0.42, 6.6, 0.55, size=32, bold=True)
    add_textbox(slide, "Implementation Stages", 2.18, 1.28, 4.2, 0.4, size=22, bold=True)

    stages = [
        ("1", "Requirement Analysis", "ATM risks and MFA requirements."),
        ("2", "System Design", "FastAPI, database, ATM and mobile screens."),
        ("3", "Biometric Enrollment", "PIN, voiceprint, and face reference."),
        ("4", "Authentication Workflow", "PIN + voice, SMS/QR link, face match."),
        ("5", "Testing & Documentation", "PIN, voice, SMS, Cloudflare, and face flow."),
    ]
    for idx, title, detail in stages:
        add_stage_line(slide, idx, title, detail, 1.95 + (int(idx) - 1) * 0.63)

    add_textbox(
        slide,
        "A systematic approach integrating PIN authentication, voice biometrics, SMS-based mobile handoff, and facial verification.",
        2.18,
        5.42,
        7.8,
        0.82,
        size=22,
        italic=True,
    )
    slide.shapes.add_picture(str(assets["integration"]), Inches(9.72), Inches(1.05), width=Inches(2.65), height=Inches(2.65))
    add_textbox(slide, "Secure System Integration", 9.85, 4.65, 2.4, 0.35, size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def add_flow_box(slide, text, x, y, w, h, fill, font_size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Times New Roman"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


def add_arrow(slide, x, y, h=0.28):
    add_textbox(slide, "↓", x, y, 0.35, h, size=20, color=BLUE, align=PP_ALIGN.CENTER)


def add_workflow_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shared_branding(slide, assets)
    add_textbox(slide, "PROPOSED SYSTEM ARCHITECTURE AND WORKFLOW", 2.18, 0.45, 8.0, 0.38, size=18, bold=True)

    cx = 5.0
    add_flow_box(slide, "User Inserts ATM Card", cx, 1.05, 3.35, 0.5, BLUE, 15)
    add_arrow(slide, 6.49, 1.58)
    add_flow_box(slide, "PIN Verification", cx, 1.85, 3.35, 0.5, BLUE, 15)
    add_arrow(slide, 6.49, 2.38)
    add_flow_box(slide, "Voice Biometric Verification", cx, 2.65, 3.35, 0.5, BLUE, 15)
    add_arrow(slide, 6.49, 3.18)

    add_flow_box(slide, "Twilio SMS Link", 3.85, 3.48, 2.45, 0.5, BLUE, 14)
    add_textbox(slide, "+", 6.46, 3.55, 0.3, 0.25, size=18, color=BLUE, align=PP_ALIGN.CENTER)
    add_flow_box(slide, "ATM QR Fallback", 6.95, 3.48, 2.45, 0.5, BLUE, 14)
    add_arrow(slide, 6.49, 4.00)

    add_flow_box(slide, "Mobile Face Verification", cx, 4.28, 3.35, 0.5, BLUE, 15)
    add_textbox(slide, "(Token expiry + face match)", cx + 0.45, 4.82, 2.5, 0.25, size=10, italic=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 6.49, 5.08)
    add_flow_box(slide, "Multi-Factor Authentication Engine", cx, 5.32, 3.35, 0.5, BLUE, 15)

    add_textbox(slide, "↓", 5.95, 5.87, 0.35, 0.25, size=16, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "↓", 7.05, 5.87, 0.35, 0.25, size=16, color=BLUE, align=PP_ALIGN.CENTER)
    add_flow_box(slide, "Access Granted", 4.05, 6.10, 1.9, 0.42, GREEN, 13)
    add_flow_box(slide, "Access Denied", 7.15, 6.10, 1.9, 0.42, RED, 13)
    add_textbox(slide, "Transaction Proceed", 4.45, 6.60, 1.1, 0.2, size=9, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Alert / Retry", 7.63, 6.60, 1.0, 0.2, size=9, align=PP_ALIGN.CENTER)


def build_pptx(page_paths: list[Path], assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Remove default empty slide if a template creates one.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for idx, path in enumerate(page_paths, start=1):
        if idx == 7:
            add_methodology_slide(prs, assets)
        elif idx == 8:
            add_workflow_slide(prs, assets)
        else:
            add_picture_slide(prs, path)
    prs.save(PPTX_PATH)


def font(size: int, bold: bool = False):
    candidates = [
        r"C:\Windows\Fonts\timesbd.ttf" if bold else r"C:\Windows\Fonts\times.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def pil_text(draw, xy, text, size=34, bold=False, italic=False, fill=(0, 0, 0), anchor=None, align="left"):
    draw.multiline_text(xy, text, font=font(size, bold=bold), fill=fill, anchor=anchor, spacing=8, align=align)


def make_preview_base(assets: dict[str, Path]) -> Image.Image:
    canvas = Image.new("RGB", (1920, 1080), "white")
    footer = Image.open(assets["footer"]).convert("RGB").resize((1920, 98))
    chevron = Image.open(assets["chevron"]).convert("RGB").resize((222, 205))
    canvas.paste(chevron, (0, 0))
    canvas.paste(footer, (0, 982))
    return canvas


def draw_box(draw, xyxy, text, fill, size=25):
    draw.rectangle(xyxy, fill=fill)
    x0, y0, x1, y1 = xyxy
    pil_text(draw, ((x0 + x1) / 2, (y0 + y1) / 2), text, size=size, bold=True, fill=(255, 255, 255), anchor="mm", align="center")


def create_previews(assets: dict[str, Path], page_paths: list[Path]) -> list[Path]:
    previews: list[Path] = []
    blue = (79, 129, 189)
    green = (20, 137, 64)
    red = (193, 23, 24)

    methodology = make_preview_base(assets)
    draw = ImageDraw.Draw(methodology)
    pil_text(draw, (313, 55), "METHODOLOGY", size=45, bold=True)
    pil_text(draw, (313, 160), "Implementation Stages", size=31, bold=True)
    lines = [
        ("1.", "Requirement Analysis", "- ATM risks and MFA requirements."),
        ("2.", "System Design", "- FastAPI, database, ATM and mobile screens."),
        ("3.", "Biometric Enrollment", "- PIN, voiceprint, and face reference."),
        ("4.", "Authentication Workflow", "- PIN + voice, SMS/QR link, face match."),
        ("5.", "Testing & Documentation", "- PIN, voice, SMS, Cloudflare, and face flow."),
    ]
    for i, (n, title, detail) in enumerate(lines):
        y = 255 + i * 83
        pil_text(draw, (313, y), n, size=33)
        pil_text(draw, (365, y), title, size=30, bold=True)
        pil_text(draw, (747, y), detail, size=28)
    pil_text(draw, (313, 780), "A systematic approach integrating PIN authentication, voice biometrics,\nSMS-based mobile handoff, and facial verification.", size=31)
    integration = Image.open(assets["integration"]).convert("RGB").resize((385, 385))
    methodology.paste(integration, (1395, 155))
    pil_text(draw, (1588, 650), "Secure System Integration", size=23, bold=True, fill=blue, anchor="mm")
    path = PREVIEW_DIR / "slide_07_methodology_preview.png"
    methodology.save(path)
    previews.append(path)

    workflow = make_preview_base(assets)
    draw = ImageDraw.Draw(workflow)
    pil_text(draw, (313, 70), "PROPOSED SYSTEM ARCHITECTURE AND WORKFLOW", size=30, bold=True)
    def down_arrow(x, y):
        pil_text(draw, (x, y), "↓", size=37, fill=blue, anchor="mm")
    draw_box(draw, (720, 150, 1200, 215), "User Inserts ATM Card", blue)
    down_arrow(960, 250)
    draw_box(draw, (720, 290, 1200, 355), "PIN Verification", blue)
    down_arrow(960, 390)
    draw_box(draw, (720, 430, 1200, 495), "Voice Biometric Verification", blue)
    down_arrow(960, 525)
    draw_box(draw, (560, 501, 895, 573), "Twilio SMS Link", blue, 23)
    pil_text(draw, (960, 536), "+", size=32, fill=blue, anchor="mm")
    draw_box(draw, (1025, 501, 1360, 573), "ATM QR Fallback", blue, 23)
    down_arrow(960, 602)
    draw_box(draw, (720, 616, 1200, 688), "Mobile Face Verification", blue)
    pil_text(draw, (960, 733), "(Token expiry + face match)", size=18, anchor="mm")
    down_arrow(960, 755)
    draw_box(draw, (720, 766, 1200, 838), "Multi-Factor Authentication Engine", blue, 24)
    draw_box(draw, (575, 878, 845, 939), "Access Granted", green, 21)
    draw_box(draw, (1075, 878, 1345, 939), "Access Denied", red, 21)
    path = PREVIEW_DIR / "slide_08_workflow_preview.png"
    workflow.save(path)
    previews.append(path)

    thumbs = []
    for idx, page_path in enumerate(page_paths, start=1):
        source = previews[0] if idx == 7 else previews[1] if idx == 8 else page_path
        img = Image.open(source).convert("RGB")
        img.thumbnail((420, 236))
        tile = Image.new("RGB", (460, 286), "white")
        tile.paste(img, ((460 - img.width) // 2, 34))
        label_draw = ImageDraw.Draw(tile)
        pil_text(label_draw, (20, 12), f"Slide {idx}", size=18, bold=True)
        thumbs.append(tile)

    cols = 4
    rows = math.ceil(len(thumbs) / cols)
    sheet = Image.new("RGB", (cols * 460, rows * 286), (245, 245, 245))
    for idx, thumb in enumerate(thumbs):
        x = (idx % cols) * 460
        y = (idx // cols) * 286
        sheet.paste(thumb, (x, y))
    contact = PREVIEW_DIR / "updated_deck_contact_sheet.png"
    sheet.save(contact)
    previews.append(contact)
    return previews


def validate() -> None:
    prs = Presentation(PPTX_PATH)
    if len(prs.slides) != 11:
        raise RuntimeError(f"Expected 11 slides, found {len(prs.slides)}")
    if PPTX_PATH.stat().st_size < 500_000:
        raise RuntimeError("Generated PPTX is unexpectedly small")


def main() -> None:
    ensure_dirs()
    page_paths = render_pdf_pages()
    assets = crop_assets(page_paths)
    build_pptx(page_paths, assets)
    previews = create_previews(assets, page_paths)
    validate()
    print(f"PPTX: {PPTX_PATH}")
    print(f"Slides: 11")
    print("Previews:")
    for preview in previews:
        print(f"- {preview}")


if __name__ == "__main__":
    main()
